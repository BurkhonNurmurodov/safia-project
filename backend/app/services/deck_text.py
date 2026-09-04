"""How much room a string needs when PowerPoint draws it.

A generated deck is laid out in absolute inches — there is no browser to
reflow it and no `overflow: hidden` to catch a mistake. A PowerPoint text box
does not clip: text longer than the box is drawn straight out of it, on top of
whatever was placed underneath. So the WEEKLY REPORT'S TEXT OVERLAPPED
(2026-09-04, reported by the operator): an event note clipped to 96 characters
took two lines inside a box built for one, and its second line landed exactly
on the «61 min · 8811 · 30-avgust» meta line drawn 0.21" below it.

The defect was never one slide. It is the shape the deck was written in —
a character cap (`_clip(s, 96)`) standing in for a WIDTH the author guessed at,
in a column whose width depends on how many cards the week happened to produce.
A guess like that is right at one size and wrong at every other, and there are
forty of them, so fixing the one slide somebody noticed would leave the rest to
be found by whoever reads next week's file.

**This module is the fix, and the rule is that the DRAWING PRIMITIVE holds the
invariant, not the call site.** `ojidaniya_deck._text()` measures every string
against the box it was handed, wraps it itself, and trims what genuinely cannot
fit — so a slide can no longer be written in a way that overlaps, whatever any
future author passes it. Same reasoning the `SegmentedToggle` overflow rule
records: an invariant every caller has to remember is one the template does not
hold.

**Estimation, and which way it is allowed to be wrong.** The advance widths
below are Calibri's own, per 1000 em, with Cyrillic — the register spells
brigadirs in two alphabets — and the punctuation this deck actually uses. They
are not exact for every renderer (PowerPoint on Windows, Google Slides
substituting metric-compatible Carlito, LibreOffice), so `SAFETY` biases every
measurement WIDE. That choice is the whole safety argument: over-estimating
breaks a line one word early, which nobody notices; under-estimating puts text
back on top of text, which is the bug. Never tune these numbers to make text
fit — widen the box instead.

**Lines are broken HERE, not by the renderer.** `wrap()` returns the lines and
the caller draws them as explicit breaks, so the number of lines measured is
the number of lines drawn. Left to PowerPoint, a width we guessed slightly
narrow would let it fit one more word per line — or one fewer, and one more
line than the box has room for. Wrapping against an over-estimated width can
only ever produce lines the real font also fits.
"""
from __future__ import annotations

# ── advance widths, per 1000 em ──────────────────────────────────────────────
# Calibri regular. The deck's body font, and the base every other font is
# scaled from.
_W: dict[str, int] = {}


def _load(table: str, widths: str) -> None:
    for ch, w in zip(table, widths.split()):
        _W[ch] = int(w)


_load(" !\"#$%&'()*+,-./",
      "226 296 400 507 507 678 634 218 302 302 439 498 248 306 248 419")
_load("0123456789", "507 507 507 507 507 507 507 507 507 507")
_load(":;<=>?@", "268 268 498 498 498 445 850")
_load("ABCDEFGHIJKLMNOPQRSTUVWXYZ",
      "579 544 533 615 488 459 631 623 252 319 520 420 855 646 662 517 673 "
      "543 459 487 642 567 890 519 487 468")
_load("[\\]^_`", "315 419 315 498 393 322")
_load("abcdefghijklmnopqrstuvwxyz",
      "479 525 423 525 498 305 471 525 229 239 455 229 799 525 527 525 525 "
      "349 391 335 525 452 715 433 453 395")
_load("{|}~", "314 460 314 498")

# Cyrillic. A supervisor's name is printed as the register spells it, and about
# a third of the register is Cyrillic — measured properly rather than folded
# into a default, because these are the strings most likely to sit in a narrow
# card heading.
_load("абвгдежзийклмнопрстуфхцчшщъыьэюя",
      "479 500 480 400 520 498 720 430 540 540 480 520 620 530 527 530 525 "
      "423 430 460 700 433 540 500 760 780 560 700 460 460 720 480")
_load("АБВГДЕЖЗИЙКЛМНОПРСТУФХЦЧШЩЪЫЬЭЮЯ",
      "579 544 544 488 615 488 890 500 646 646 620 615 855 623 662 517 543 "
      "533 487 567 700 519 646 543 890 890 700 890 544 533 890 567")
_load("ЎўҚқҒғҲҳЁёЪъ", "642 525 620 480 488 400 623 525 488 498 700 560")

# The marks this deck sets, and nothing else — a glyph nobody prints does not
# belong in a table somebody has to trust.
_load("«»·–—…‚„“”‘’≈≠≤≥±−×÷→←▲▼●■□✓№°",
      "427 427 305 498 900 808 248 400 400 400 218 218 498 498 498 498 498 "
      "498 498 498 838 838 700 700 600 600 600 600 800 400")

# What an unmeasured character costs. Split by case because the two are far
# apart, and set ABOVE the average of each class on purpose: an unknown glyph
# must never be cheaper than it draws.
_DEFAULT_UPPER = 700
_DEFAULT_LOWER = 560

# Cambria is the heading face and is wider than Calibri, most of all in
# capitals. One scale rather than a second table: headings are short, and the
# cost of a scale that is slightly generous is a heading trimmed a character
# early.
_FONT_SCALE = {"Calibri": 1.0, "Cambria": 1.12}
_BOLD_SCALE = 1.03

# Everything measured here is multiplied by this. See the module docstring:
# the direction of the error is the safety argument.
SAFETY = 1.03

# One line of text, as a multiple of the font size — the font's own ascent +
# descent + line gap. PowerPoint's percentage line spacing multiplies THIS.
_LINE = {"Calibri": 1.22, "Cambria": 1.17}
_LINE_DEFAULT = 1.22

ELLIPSIS = "…"


def char_w(ch: str) -> int:
    w = _W.get(ch)
    if w is not None:
        return w
    return _DEFAULT_UPPER if ch.isupper() else _DEFAULT_LOWER


def width_pt(s: str, size: float, font: str = "Calibri", bold: bool = False) -> float:
    """How wide `s` draws, in points, biased wide."""
    if not s:
        return 0.0
    units = sum(char_w(ch) for ch in s)
    scale = _FONT_SCALE.get(font, 1.12) * (_BOLD_SCALE if bold else 1.0)
    return units / 1000.0 * size * scale * SAFETY


def line_h_pt(size: float, font: str = "Calibri", spacing: float | None = None) -> float:
    """One line's height in points, honouring an explicit line-spacing multiple."""
    return size * _LINE.get(font, _LINE_DEFAULT) * (spacing or 1.0)


def block_h_in(n_lines: int, size: float, font: str = "Calibri",
               spacing: float | None = None) -> float:
    """The height in INCHES `n_lines` occupy — what a caller stacking items
    must advance by."""
    return max(0, n_lines) * line_h_pt(size, font, spacing) / 72.0


def lines_that_fit(height_in: float, size: float, font: str = "Calibri",
                   spacing: float | None = None) -> int:
    """How many lines a box that tall can show. At least one: every box in the
    deck is sized for a line, and a zero would print nothing at all."""
    lh = line_h_pt(size, font, spacing)
    if lh <= 0:
        return 1
    # A half-point of slack, because the boxes were hand-sized in hundredths of
    # an inch and a box built for exactly two lines must not read as one.
    return max(1, int((height_in * 72.0 + 0.5) / lh))


def _break_word(word: str, avail_pt: float, size: float, font: str,
                bold: bool) -> list[str]:
    """A single word too long for the line — a code, a URL, a leader typing
    without spaces. Broken by character, because the alternative is one line
    running off the card."""
    out, cur = [], ""
    for ch in word:
        trial = cur + ch
        if cur and width_pt(trial, size, font, bold) > avail_pt:
            out.append(cur)
            cur = ch
        else:
            cur = trial
    if cur:
        out.append(cur)
    return out or [word]


def wrap(text: str, width_in: float, size: float, font: str = "Calibri",
         bold: bool = False) -> list[str]:
    """`text` broken into lines that each fit `width_in`. Explicit newlines are
    honoured as breaks the author asked for."""
    avail = max(width_in, 0.05) * 72.0
    lines: list[str] = []
    for hard in (text or "").split("\n"):
        words = hard.split()
        if not words:
            lines.append("")
            continue
        cur = ""
        for word in words:
            trial = f"{cur} {word}" if cur else word
            if cur and width_pt(trial, size, font, bold) > avail:
                lines.append(cur)
                cur = word
            else:
                cur = trial
            if width_pt(cur, size, font, bold) > avail:
                pieces = _break_word(cur, avail, size, font, bold)
                lines.extend(pieces[:-1])
                cur = pieces[-1]
        lines.append(cur)
    return lines


def fit(text: str, width_in: float, size: float, font: str = "Calibri",
        bold: bool = False, max_lines: int = 1) -> tuple[list[str], bool]:
    """The lines to draw, and whether anything had to be dropped.

    The last line keeps as many whole words as fit beside an ellipsis, so a
    trimmed sentence still ends on a word rather than mid-syllable.
    """
    lines = wrap(text, width_in, size, font, bold)
    if len(lines) <= max_lines:
        return lines, False
    avail = max(width_in, 0.05) * 72.0
    kept = lines[: max(1, max_lines)]
    words = kept[-1].split()
    while words:
        trial = " ".join(words).rstrip(" ,.;:·-—") + ELLIPSIS
        if width_pt(trial, size, font, bold) <= avail:
            kept[-1] = trial
            return kept, True
        words.pop()
    kept[-1] = ELLIPSIS
    return kept, True


# ── the invariant, checked on the finished file ─────────────────────────────
def check_layout(prs, page_w_in: float, page_h_in: float,
                 eps: float = 0.006) -> list[str]:
    """Every way a generated deck can put text where it does not belong.

    Two rules, and it is the PAIR that makes overlapping text impossible:

    1. **No box holds more lines than it has room for.** `_text` guarantees
       this per box — it wraps and trims against the box it was handed — so a
       string can never be drawn outside its own rectangle.
    2. **No two text boxes intersect, and none leaves the page.** This is what
       carries rule 1 from a box to a SLIDE. Rule 1 alone only promises that
       text stays inside its rectangle; if two rectangles overlap, two strings
       that each obey rule 1 still land on top of each other.

    Rule 1 cannot be checked here — the file records boxes, not glyphs — and
    does not need to be: it holds by construction. Rule 2 is geometry, so it
    is checked on the finished presentation, on real data, every time one is
    built. It returns findings rather than raising: a report that arrives with
    a cosmetic flaw is worth more than no report, the same call
    `deck_narrative` makes about a Gemini failure.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE      # local: this is a checker

    out: list[str] = []
    for n, slide in enumerate(prs.slides, start=1):
        boxes = []
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            text = (sh.text_frame.text or "").strip()
            if not text:
                continue
            x0, y0 = sh.left / 914400, sh.top / 914400
            boxes.append((x0, y0, x0 + sh.width / 914400,
                          y0 + sh.height / 914400, " ".join(text.split())[:60]))
        for x0, y0, x1, y1, t in boxes:
            if x0 < -eps or y0 < -eps or x1 > page_w_in + eps or y1 > page_h_in + eps:
                out.append(f"slide {n}: «{t}» is off the page")
        for i in range(len(boxes)):
            a = boxes[i]
            for j in range(i + 1, len(boxes)):
                b = boxes[j]
                if (a[0] < b[2] - eps and b[0] < a[2] - eps
                        and a[1] < b[3] - eps and b[1] < a[3] - eps):
                    out.append(f"slide {n}: «{a[4]}» overlaps «{b[4]}»")
    return out
