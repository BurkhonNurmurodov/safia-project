"""The weekly Ojidaniya report as a PowerPoint deck.

Once a week somebody has to say what the plant waited on, who it happened to
and what to do about it. That was being done by hand — read every leader's
entry, add the minutes up per category and per brigadir, find the three that
matter, write the root causes, lay out fourteen slides. This does the whole of
it from the register, for the period `services/report_week.py` defines.

The rules it is built on, each of which is a decision somebody made:

**Every figure comes from the page's own computation.** `_downtime()` and
`_cell_detail()` in `routers/downtime.py` are what `/downtime` itself reads;
this module is handed their output and formats it. It never queries and never
re-derives, so the deck and the screen can never state different numbers — the
same rule `ojidaniya_export.py` follows for the workbook.

**Two totals, and both are named.** A unit's day is the headcount-weighted MEAN
of its cells (`Σ(N·T)÷ΣN`), so the events listed under a brigadir add up to
something else — usually much more. The KPI slides carry the weighted figure
that `/downtime` charts; the per-cell and per-event slides carry the cells' own
sums and SAY they are the cells' own sums. Printing one of them everywhere
would either contradict the page or leave an unexplained gap.

**The deck ignores the page's filters.** It is a fixed weekly report about one
plant, both shifts, every supervisor, headline on the stopped half with the
not-stopped half kept as an aside — the shape of the hand-made deck it
replaces. What is on screen when the button is pressed changes nothing, which
is why the button confirms with the scope written out.

**Its category scope is the ЗАГРУЗКА's** (the operator's ruling, 2026-09-04):
`sheets_reader.OJIDANIYA_ONLY_CATS` — «Cat H», Тозалаш, the cleaning the shift
does rather than a stoppage it suffers — is not in this file at all. The
filtering happens where the data is fetched, `DECK_KPI_ONLY` in
`routers/downtime.py`, so this module never sees a dropped category and has
nothing to subtract. Do not add a second filter here.

**No text on a slide may ever land on other text.** Everything here is
absolutely positioned and a PowerPoint text box does not clip, so a string
longer than its box is drawn over whatever is underneath — which is what
happened to the event notes on 2026-09-04. The rule is held in TWO places and
neither is a call site: `_text` wraps and trims every string against the box it
was handed (`services/deck_text.py`), and `deck_text.check_layout` verifies on
the finished file that no two boxes intersect. A slide author states a box; the
box is the promise.

**It is always Uzbek**, whoever presses it, so the words live here rather than
travelling from the client as they do for the workbook. The category labels are
a copy of the `downtime.cat.*` uz bundle in `frontend/src/i18n/translations.js`
and the colours mirror `catColor()` — keep them in step.

**Leaders' notes are quoted, never rewritten.** A note is the recorded evidence
of a shift. It is transliterated to Latin where it was typed in Cyrillic (the
register holds both, in the same column) and otherwise reproduced as written,
typos included. `services/deck_narrative.py` writes the commentary AROUND the
notes; it is not allowed to restate them.
"""
from __future__ import annotations

import io
import logging
from collections import defaultdict
from datetime import date, timedelta

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.text.text import Font as _Font
from pptx.util import Emu, Inches, Pt

from app.services import deck_text
from app.translit import transliterate

log = logging.getLogger(__name__)


# ── the deck's own look ──────────────────────────────────────────────────────
# Lifted from the hand-made deck this replaces, so a reader who has seen last
# week's file recognises this one. Warm browns and cream, not the app's chrome.
INK    = RGBColor(0x24, 0x1B, 0x14)   # headings, the numbers that matter
INK2   = RGBColor(0x3A, 0x2B, 0x1E)   # body text
MUTED  = RGBColor(0x6A, 0x5D, 0x52)   # secondary text
FAINT  = RGBColor(0x97, 0x89, 0x7C)   # captions, footers
CREAM  = RGBColor(0xF7, 0xF3, 0xEE)   # page
LINE   = RGBColor(0xE7, 0xE0, 0xD8)   # borders
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
INNER  = RGBColor(0xFA, 0xF7, 0xF3)   # recessed panel inside a card
GOLD   = RGBColor(0xC9, 0x9E, 0x63)   # brand accent
GOLD2  = RGBColor(0xE8, 0xB0, 0x4B)   # accent on the dark cover
BROWN  = RGBColor(0x7A, 0x4A, 0x21)
DARK   = RGBColor(0x2A, 0x1B, 0x12)   # cover ground
ONDARK = RGBColor(0xB9, 0xA8, 0x97)
RED    = RGBColor(0x9F, 0x33, 0x32)
REDBG  = RGBColor(0xFB, 0xE6, 0xE5)
GREEN  = RGBColor(0x2F, 0x6B, 0x34)
GREENBG = RGBColor(0xDB, 0xEE, 0xDB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

W, H = 13.333, 7.5          # inches, 16:9
M = 0.55                    # page margin
CW = W - 2 * M              # content width

# Ojidaniya categories in canonical A→Z order — MUST mirror
# `routers/idle_cell.IDLE_CATEGORIES` and `frontend/src/components/idle/
# categories.js`. The colour is that list's position into CATEGORY_COLORS, the
# same rule `catColor()` applies, so a category is one hue across the platform
# and this file.
CATS: list[tuple[str, str, str, str]] = [
    # (stored name, code, short label, full label)
    ("Cat A",  "A",  "Xoladilnikdan mahsulot",  "Xoladilnikdan mahsulot kutish"),
    ("Cat A2", "A2", "Xoladilnikdan GP",        "Xoladilnikdan GP mahsulot kutish"),
    ("Cat B",  "B",  "Uskuna nosozligi",        "Oborudivaniya buzilishi"),
    ("Cat C",  "C",  "List / vagonetka",        "List/vaganetka kutish"),
    ("Cat D",  "D",  "Skladdan xomashyo",       "Skladdan mahsulot yoki hom ashyo kutish"),
    ("Cat D2", "D2", "Qo'shimcha zayavka",      "Skladdan qo'shimcha zayavka orqali hom ashyo kutish"),
    ("Cat D3", "D3", "Otdellararo mahsulot",    "Otdellardan mahsulot kutish"),
    ("Cat E",  "E",  "Ichki logistika",         "Ichki logistikadan mahsulot yoki hom ashyo kutish"),
    ("Cat F",  "F",  "Texnolog qarori",         "Texnologlar qarorini kutish"),
    ("Cat G",  "G",  "Plan bo'limi",            "Plan bo'limi"),
    ("Cat H",  "H",  "Tozalash",                "Tozalash"),
    ("Cat I",  "I",  "Smena topshiruvi",        "Oldingi smena ishi tugashini kutish"),
]
_CAT_BY_NAME = {c[0]: c for c in CATS}

# utils/chartPalette.js CATEGORY_COLORS, generic-first.
_PALETTE = [
    (0xEF, 0x44, 0x44), (0x22, 0xC5, 0x5E), (0x3B, 0x82, 0xF6), (0xEA, 0xB3, 0x08),
    (0xF9, 0x73, 0x16), (0xA8, 0x55, 0xF7), (0x14, 0xB8, 0xA6), (0xEC, 0x48, 0x99),
    (0x63, 0x66, 0xF1), (0x84, 0xCC, 0x16), (0x06, 0xB6, 0xD4), (0xD9, 0x46, 0xEF),
]

# Categories that show on the Ojidaniya page but never enter the загрузка —
# `sheets_reader.OJIDANIYA_ONLY_CATS`, restated here only so a slide can name
# the scope it is describing.
#
# **They do not reach this module at all** (the operator's ruling, 2026-09-04).
# The filtering happens where the deck's data is fetched — `DECK_KPI_ONLY` in
# `routers/downtime.py`, i.e. the platform's own `kpi_only` door — so `collect`
# is handed the загрузка scope already and there is nothing here to subtract or
# hide. Never add a second filter in this file: two spellings of the rule is
# how the deck and the page start disagreeing about one week.
#
# `kpi` therefore rides on every category row and is true for all of them
# today. It stays because it is a true statement about a category rather than
# about this week's data, and because lifting the ruling must not mean
# rediscovering which categories the flag was about.
OJIDANIYA_ONLY = {"Cat H"}


def cat_meta(name: str) -> tuple[str, str, str, RGBColor]:
    """(code, short label, full label, colour) for a stored category name."""
    row = _CAT_BY_NAME.get(name)
    if not row:
        return (name.replace("Cat ", ""), name, name, RGBColor(*_PALETTE[0]))
    idx = CATS.index(row)
    return (row[1], row[2], row[3], RGBColor(*_PALETTE[idx % len(_PALETTE)]))


# ── formatting ───────────────────────────────────────────────────────────────
def num(v: float) -> str:
    """1542 → «1 542». A space, the Uzbek thousands separator."""
    return f"{int(round(v)):,}".replace(",", " ")


def hours(mins: float) -> str:
    """Minutes as hours with one decimal, comma as the decimal mark."""
    return f"{mins / 60:.1f}".replace(".", ",")


def pct(v: float, digits: int = 0) -> str:
    s = f"{v:.{digits}f}".replace(".", ",")
    return f"{s}%"


def signed_pct(v: float) -> str:
    """«+12%» / «−65%». A typographic minus (U+2212), not a hyphen: at the
    display sizes the cover uses, a hyphen reads as punctuation rather than a
    sign, which is why the hand-made deck set it this way too."""
    return ("+" if v > 0 else "") + pct(v).replace("-", "−")


_UZ_MONTH = ["", "yanvar", "fevral", "mart", "aprel", "may", "iyun",
             "iyul", "avgust", "sentabr", "oktabr", "noyabr", "dekabr"]
_UZ_WDAY = ["dushanba", "seshanba", "chorshanba", "payshanba",
            "juma", "shanba", "yakshanba"]


def day_label(d: date) -> str:
    return f"{d.day}-{_UZ_MONTH[d.month]}"


def day_full(d: date) -> str:
    return f"{d.day}-{_UZ_MONTH[d.month]}, {_UZ_WDAY[d.weekday()]}"


def period_words(win: tuple[date, date]) -> str:
    """«26-avgust — 2-sentabr, 2026» — how the deck names its own period."""
    a, b = win
    if a.year == b.year:
        return f"{a.day}-{_UZ_MONTH[a.month]} — {b.day}-{_UZ_MONTH[b.month]}, {b.year}"
    return f"{a.day}-{_UZ_MONTH[a.month]}, {a.year} — {b.day}-{_UZ_MONTH[b.month]}, {b.year}"


def latin(s: str) -> str:
    """A leader's note in Latin script, otherwise exactly as they typed it.

    The register holds both alphabets in one column — roughly a third of a
    week's notes arrive in Cyrillic — and a deck that mixes them reads as two
    documents stapled together. This changes the SCRIPT and nothing else: no
    spelling is corrected, no wording touched.
    """
    return transliterate((s or "").strip(), "uz")


def _clip(s: str, n: int) -> str:
    s = " ".join((s or "").split())
    return s if len(s) <= n else s[: n - 1].rstrip() + "…"


def _one(text: str, width: float, size: float, *, font=BODY_FONT,
         bold=False) -> str:
    """`text` cut to a SINGLE line at `width`.

    For a string that shares a box with a line below it inside the same
    paragraph: `_text` trims a paragraph from the END, so a first segment that
    ran to two lines would push the second one out of the box entirely and the
    reader would lose the line nobody was worried about.
    """
    return deck_text.fit(text, width, size, font, bold, 1)[0][0]


def short_name(name: str, width: float, size: float, *, font=BODY_FONT,
               bold=False) -> str:
    """A person's name, shortened only as far as it has to be.

    A brigadir's name is the identity of the card it heads, so «Абдурахмонова
    Гулнора Шухратовна» must not become «Абдурахмонова Гулнора Шухрат…»: an
    ellipsis says the platform ran out of room, initials say who this is. The
    full name is used wherever it fits, then trailing parts become initials one
    at a time — the most informative form that fits wins — and a name still too
    wide is left to `_text`, which can always trim it.
    """
    parts = (name or "").split()
    if not parts:
        return ""
    avail = width * 72.0
    if deck_text.width_pt(name, size, font, bold) <= avail:
        return name
    for keep in range(len(parts) - 1, 0, -1):
        trial = " ".join(parts[:keep] + [f"{p[0]}." for p in parts[keep:]])
        if deck_text.width_pt(trial, size, font, bold) <= avail:
            return trial
    return name


def _sentence(s: str) -> str:
    """Leaders type in caps, or with no capital at all. Only the first letter
    is touched — a word already carrying a lowercase letter is left alone, the
    `utils/textCase.js` rule."""
    s = " ".join((s or "").split())
    if not s:
        return s
    words = s.split(" ")
    fixed = [w.capitalize() if w.isupper() and len(w) > 1 else w for w in words]
    out = " ".join(fixed)
    return out[0].upper() + out[1:] if out else out


# ── what the deck is made of ─────────────────────────────────────────────────
# A note this short says nothing a reader can act on («Uvuvuvi», «Un yoqligi»).
# They are counted for the appendix rather than dropped: a register that is
# thin in places is a fact about the week, and hiding it makes the other notes
# look more complete than they are.
THIN_NOTE_CHARS = 12

# The unit-day threshold `/downtime` flags in red. Restated so the deck's own
# tables can mark the same days the page marks.
FLAG_MIN = 50


def collect(*, cur: dict, prev: dict, events: list[dict], cell_days: list[dict],
            win: tuple[date, date], prev_win: tuple[date, date],
            factory_name: str, supervisors: list[dict]) -> dict:
    """Shape the page's own output into everything the fourteen slides need.

    `cur` / `prev` are `_downtime()` results — the weighted unit figures the
    charts on `/downtime` are drawn from. `events` and `cell_days` come from
    `_cell_detail()` — the cells' own filings, which is where every note, cell
    ranking and per-brigadir breakdown comes from. The two are kept apart on
    purpose and the slides say which one they are showing.
    """
    d_from, d_to = win
    days = [d_from + timedelta(days=i) for i in range((d_to - d_from).days + 1)]

    # ── the weighted figures: totals, categories, daily, supervisors ─────────
    def totals(res: dict) -> tuple[float, float, dict[str, float]]:
        tot = tot_ns = 0.0
        by_cat: dict[str, float] = defaultdict(float)
        for r in res["rows"]:
            tot += float(r["total"] or 0)
            tot_ns += float(r["total_ns"] or 0)
            for c, v in (r["by_category"] or {}).items():
                by_cat[c] += float(v or 0)
        return tot, tot_ns, dict(by_cat)

    total, total_ns, by_cat = totals(cur)
    p_total, p_total_ns, p_by_cat = totals(prev)

    daily_map: dict[str, float] = defaultdict(float)
    for r in cur["rows"]:
        daily_map[r["date"]] += float(r["total"] or 0)
    daily = [{"date": d, "label": day_label(d), "full": day_full(d),
              "minutes": daily_map.get(d.strftime("%d.%m.%Y"), 0.0)}
             for d in days]

    # ── the cells' own filings ──────────────────────────────────────────────
    ev_on = [e for e in events if e["stopped"]]
    ev_off = [e for e in events if not e["stopped"]]

    ev_by_cat: dict[str, list[dict]] = defaultdict(list)
    for e in ev_on:
        ev_by_cat[e["category"]].append(e)
    ev_by_sup: dict[str, list[dict]] = defaultdict(list)
    for e in ev_on:
        ev_by_sup[e["supervisor"]].append(e)

    # A cell's week is the sum of its DAILY unions — never the sum of its
    # events, which double-counts two causes overlapping on one clock.
    cell_min: dict[str, float] = defaultdict(float)
    cell_meta: dict[str, dict] = {}
    for cd in cell_days:
        if not cd["stopped"]:
            continue
        cell_min[cd["cell"]] += float(cd["union_minutes"] or 0)
        cell_meta.setdefault(cd["cell"], {"leader": cd.get("leader"),
                                          "supervisor": cd["supervisor"]})
    cell_events: dict[str, int] = defaultdict(int)
    for e in ev_on:
        cell_events[e["cell"]] += 1
    cells = sorted(
        ({"code": c, "minutes": m, "events": cell_events.get(c, 0),
          **cell_meta.get(c, {"leader": None, "supervisor": ""})}
         for c, m in cell_min.items()),
        key=lambda x: -x["minutes"])

    # ── categories, ranked, with last week beside them ──────────────────────
    cat_rows = []
    for name in sorted(set(by_cat) | set(p_by_cat) | {e["category"] for e in ev_on}):
        code, short, full, colour = cat_meta(name)
        mins = by_cat.get(name, 0.0)
        cat_rows.append({
            "name": name, "code": code, "label": short, "full": full,
            "colour": colour, "minutes": mins, "prev": p_by_cat.get(name, 0.0),
            "events": len(ev_by_cat.get(name, [])),
            "share": (mins / total * 100) if total else 0.0,
            "kpi": name not in OJIDANIYA_ONLY,
        })
    cat_rows.sort(key=lambda c: -c["minutes"])

    # ── supervisors ─────────────────────────────────────────────────────────
    p_sup = {s["manager_name"]: float(s["total"] or 0) for s in prev["summary"]}
    sups = []
    for s in cur["summary"]:
        n = s["manager_name"]
        mins = float(s["total"] or 0)
        own = ev_by_sup.get(n, [])
        top = defaultdict(float)
        for e in own:
            top[e["category"]] += e["minutes"]
        sups.append({
            "name": n, "shift": s.get("shift"), "minutes": mins,
            "prev": p_sup.get(n, 0.0),
            "flagged_days": int(s.get("flagged_days") or 0),
            "share": (mins / total * 100) if total else 0.0,
            "events": len(own),
            "top_cat": max(top, key=top.get) if top else None,
            # The cells' own sum for this unit — deliberately a different
            # number from `minutes`, and labelled as such wherever both appear.
            "cell_minutes": sum(e["minutes"] for e in own),
        })
    sups = [s for s in sups if s["minutes"] > 0 or s["events"]]
    sups.sort(key=lambda s: -s["minutes"])

    # ── which days answered from where, and how thin the register is ────────
    # WHERE a day's minutes came from — the cells' interval model or the
    # «Смена отчёт» row — is deliberately carried onto NO slide, and that is a
    # ruling (the operator, 2026-09-03), not an oversight to be tidied up: the
    # deck must not mention that the measurement changed. Do not re-add a
    # source note, a comparability flag or a caveat without asking.
    #
    # Consequence to know before touching the comparison slides: across
    # `idle_source.CELLS_FROM` the week-over-week percentage sets two DIFFERENT
    # measurements against each other, so it can show a large improvement that
    # nobody caused — the first real run read −64.5% with the entire comparison
    # week taken off the sheet. The deck states it as an improvement anyway.
    # That is the specified behaviour.
    filed = {e["supervisor"] for e in events}
    silent = sorted(s["name"] for s in supervisors if s["name"] not in filed)
    thin = [e for e in ev_on if len(latin(e["note"])) < THIN_NOTE_CHARS]

    delta = ((total - p_total) / p_total * 100) if p_total else 0.0

    return {
        "factory": factory_name,
        "window": win, "prev_window": prev_win,
        "period": period_words(win), "prev_period": period_words(prev_win),
        "days": days, "daily": daily,
        "total": total, "prev_total": p_total, "delta_pct": delta,
        "total_ns": total_ns, "prev_total_ns": p_total_ns,
        "events": len(ev_on), "events_ns": len(ev_off),
        "avg_event": (sum(e["minutes"] for e in ev_on) / len(ev_on)) if ev_on else 0.0,
        "categories": cat_rows,
        "supervisors": sups,
        "sup_count": len(sups),
        "cells": cells,
        "cell_count": len(cells),
        "events_all": events,
        "events_on": ev_on, "events_off": ev_off,
        "by_cat_events": {k: v for k, v in ev_by_cat.items()},
        "by_sup_events": {k: v for k, v in ev_by_sup.items()},
        "quality": {
            "silent_supervisors": silent,
            "thin_notes": len(thin),
            "cyrillic_notes": sum(1 for e in ev_on
                                  if latin(e["note"]) != (e["note"] or "").strip()),
            "flag_days": sum(1 for r in cur["rows"] if r.get("flagged")),
        },
    }


# ── drawing primitives ───────────────────────────────────────────────────────
from pptx.enum.shapes import MSO_SHAPE                       # noqa: E402


def _blank(prs: Presentation):
    """A slide with nothing on it. Every slide here is absolutely positioned —
    there is one layout and it carries no placeholders, exactly as the deck
    this replaces was built."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line=None, radius=None, lw=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        # `adjustments` is a fraction of the SHORT side, so a fixed radius in
        # inches keeps corners identical on a tall card and a flat chip.
        shape.adjustments[0] = min(0.5, radius / min(w, h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    if shape.has_text_frame:
        shape.text_frame.clear()
    return shape


def _tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


# How many strings this build had to trim. A trim is not a failure — it is the
# guarantee working — but it is the one thing worth knowing about a finished
# file, so `build()` logs the tally and names each trimmed string at DEBUG.
# Silence here means every slide showed everything it was given.
_TRIMS: list[str] = []


def _p(tf, text, *, width, size=11, color=INK2, bold=False, font=BODY_FONT,
       align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=None,
       first=False, caps=False, max_lines=1):
    """One paragraph, broken into lines HERE rather than by the renderer.

    `width` and `max_lines` come from the BOX, never from the caller's
    judgement — see `_text`. Every line is its own run separated by an explicit
    break, so the number of lines measured is the number of lines drawn and a
    paragraph can never grow taller than the box it was given.

    `first` reuses the frame's own empty paragraph, which python-pptx always
    creates — appending without it leaves a blank line at the top of every text
    box.
    """
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    para.space_after = Pt(space_after)
    if line:
        para.line_spacing = line
    # Capitals are WIDER, so the case change happens before the measurement and
    # not after the run was written, as it used to.
    body = (text or "").upper() if caps else (text or "")
    lines, trimmed = deck_text.fit(body, width, size, font, bold, max_lines)
    if trimmed:
        _TRIMS.append(body)
        log.debug("deck: trimmed to %d line(s) in %.2f in: %s", max_lines, width, body)

    for i, ln in enumerate(lines):
        if i:
            # The break carries the same size as the text around it: an
            # <a:br/> with no properties of its own is laid out at the theme's
            # default 18pt, which would make a two-line block taller than the
            # two lines that were measured.
            para.add_line_break()
            brk = _Font(para._p[-1].get_or_add_rPr())
            brk.size = Pt(size)
            brk.name = font
        run = para.add_run()
        run.text = ln
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.name = font
        f.color.rgb = color
    return para


def _text(slide, x, y, w, h, text, **kw):
    """The common case: one paragraph in its own box — and THE place the deck's
    no-overflow rule is enforced.

    A PowerPoint text box does not clip, so a string longer than its box is
    drawn straight over whatever sits below it. Every call already states a
    real height; this works out from that height how many lines the box can
    show, wraps the text to the box's own width and trims the remainder. A
    caller cannot opt out and does not have to remember to. `max_lines` is
    accepted only so a block measured elsewhere (the event stack) can say what
    it measured, and it is still capped by what the box holds.
    """
    anchor = kw.pop("anchor", MSO_ANCHOR.TOP)
    size = kw.get("size", 11)
    font = kw.get("font", BODY_FONT)
    room = deck_text.lines_that_fit(h, size, font, kw.get("line"))
    asked = kw.pop("max_lines", None)
    tf = _tb(slide, x, y, w, h, anchor)
    _p(tf, text, first=True, width=w,
       max_lines=min(room, asked) if asked else room, **kw)
    return tf


def _chrome(slide, eyebrow: str, title: str, page: int, footer: str,
            *, dark=False):
    """Every slide wears the same three marks: a gold eyebrow saying what kind
    of slide this is, the question it answers, and a footer naming the report
    and the page. The eyebrow is what lets a reader flick through fourteen
    slides and find the one they want."""
    ink = WHITE if dark else INK
    sub = ONDARK if dark else FAINT
    _text(slide, M, 0.42, CW, 0.22, eyebrow, size=9.5, color=GOLD2 if dark else GOLD,
          bold=True, caps=True)
    _text(slide, M, 0.66, CW, 0.62, title, size=25, color=ink, font=HEAD_FONT, bold=True)
    _text(slide, M, H - 0.46, CW - 0.5, 0.22, footer, size=8.5, color=sub)
    _text(slide, W - M - 0.5, H - 0.46, 0.5, 0.22, str(page), size=8.5, color=sub,
          align=PP_ALIGN.RIGHT)


def _card(slide, x, y, w, h, *, fill=CARD, line=LINE, radius=0.12):
    return _rect(slide, x, y, w, h, fill=fill, line=line, radius=radius)


def _chip(slide, x, y, w, h, text, *, fill, color, size=9):
    _rect(slide, x, y, w, h, fill=fill, radius=0.06)
    _text(slide, x, y + (h - 0.17) / 2, w, 0.2, text, size=size, color=color,
          bold=True, align=PP_ALIGN.CENTER)


def _bar_chart(slide, x, y, w, h, categories, values, colours, *,
               horizontal=True, number_format='#,##0', gap=60):
    """A native chart, not a picture — the reader can click into it, and it
    stays sharp at any zoom. One series; the colours are per POINT, because
    every bar here is a different category, unit or cell."""
    data = CategoryChartData()
    data.categories = categories
    data.add_series("min", values)
    kind = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    frame = slide.shapes.add_chart(kind, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = frame.chart
    chart.has_legend = False
    chart.has_title = False

    plot = chart.plots[0]
    plot.gap_width = gap
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.number_format = number_format
    labels.number_format_is_linked = False
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    labels.font.size = Pt(8.5)
    labels.font.name = BODY_FONT
    labels.font.color.rgb = MUTED

    series = plot.series[0]
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colours[i % len(colours)]
        point.format.line.fill.background()

    cat_ax, val_ax = chart.category_axis, chart.value_axis
    cat_ax.has_major_gridlines = False
    cat_ax.tick_labels.font.size = Pt(9)
    cat_ax.tick_labels.font.name = BODY_FONT
    cat_ax.tick_labels.font.color.rgb = INK2
    cat_ax.format.line.color.rgb = LINE
    val_ax.has_major_gridlines = False
    val_ax.visible = False
    return chart


# ── the prose slots, and what stands in for them ─────────────────────────────
SHIFT_MIN = 480     # the flat shift base the загрузка uses everywhere

NO_AI = "AI izohi hozircha mavjud emas — raqamlar, jadvallar va grafiklar to'liq."


def _ai(narr: dict | None, key: str, default=None):
    """A prose slot, or the stand-in. Never raises on a partial answer: the
    model may fill eleven fields and drop one, and eleven good slides plus one
    honest gap beats no deck at all."""
    if not narr:
        return default
    v = narr.get(key)
    if isinstance(v, str):
        v = v.strip()
    return v or default


def _ai_list(narr: dict | None, key: str, n: int) -> list[dict]:
    v = _ai(narr, key) or []
    return v[:n] if isinstance(v, list) else []


# ── slide 1 · cover ──────────────────────────────────────────────────────────
def _cover(prs, d: dict, narr):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=DARK)
    _rect(s, 0, 0, W, 0.09, fill=GOLD)

    _text(s, M, 1.05, CW, 0.25, "HAFTALIK OPERATSION TAHLIL · ISHLAB CHIQARISH",
          size=10, color=GOLD2, bold=True)
    _text(s, M, 1.42, 8.6, 1.5, "Yacheykalardagi kutish\nvaqtlari tahlili",
          size=40, color=WHITE, font=HEAD_FONT, bold=True, line=1.05)
    _text(s, M, 3.05, 8.6, 0.3,
          f"{d['sup_count']} brigadir jurnali  ·  {d['events']} hodisa  ·  "
          f"{d['period']}", size=13, color=ONDARK)
    _text(s, M, 3.42, 8.6, 0.28, d["factory"], size=11, color=GOLD2, bold=True)

    # The headline number, in hours — 26 soat lands where 1 542 daqiqa does not.
    _text(s, M, 4.05, 5.0, 1.0, f"{hours(d['total'])} soat",
          size=52, color=GOLD2, font=HEAD_FONT, bold=True)
    shifts = f"{d['total'] / SHIFT_MIN:.1f}".replace(".", ",")
    move = (f" · o'tgan haftadan {signed_pct(d['delta_pct'])}"
            if d["prev_total"] else " · taqqoslash uchun ma'lumot yo'q")
    _text(s, M, 5.05, 7.4, 0.5,
          f"jami kutish — {num(d['total'])} daqiqa (≈ {shifts} smena){move}",
          size=11.5, color=ONDARK)

    # Every category the week actually had, as chips in canonical order, so the
    # cover already says what the report is made of.
    present = [c for c in d["categories"] if c["minutes"] > 0 or c["events"]]
    order = {name: i for i, (name, *_r) in enumerate(CATS)}
    present.sort(key=lambda c: order.get(c["name"], 99))
    x, y, cw = 8.35, 4.05, 0.52
    for i, c in enumerate(present[:12]):
        col, row = i % 4, i // 4
        _chip(s, x + col * (cw + 0.12), y + row * 0.5, cw, 0.36, c["code"],
              fill=c["colour"], color=WHITE, size=11)

    top3 = " · ".join(c["code"] for c in d["categories"][:3])
    _text(s, 8.35, y + 3 * 0.5 + 0.18, 4.4, 0.3,
          f"eng katta uchtasi: {top3} — 4- va 7-sahifada", size=9.5, color=ONDARK)

    _text(s, M, H - 0.5, CW, 0.24,
          "Safia · Ishlab chiqarish boshqaruvi paneli", size=9, color=RGBColor(0x7A, 0x6A, 0x5C))


# ── slide 2 · executive summary ──────────────────────────────────────────────
def _summary(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "Qisqacha", "Asosiy xulosalar", page, footer)

    shifts = f"{d['total'] / SHIFT_MIN:.1f}".replace(".", ",")
    move = signed_pct(d["delta_pct"]) if d["prev_total"] else "—"
    kpis = [
        (f"{num(d['total'])} min", "jami kutish vaqti",
         f"≈ {hours(d['total'])} soat = {shifts} smena ({move})"),
        (num(d["events"]), "yacheykani to'xtatgan hodisa",
         f"+{d['events_ns']} to'xtatmagan ({num(d['total_ns'])} min)"),
        (f"{num(d['avg_event'])} min", "o'rtacha bitta hodisa",
         f"eng uzuni {num(max((e['minutes'] for e in d['events_on']), default=0))} min"),
        (num(d["sup_count"]), "brigadir jurnal berdi",
         f"{d['cell_count']} yacheykada kutish qayd etildi"),
    ]
    cw = (CW - 3 * 0.16) / 4
    for i, (big, label, sub) in enumerate(kpis):
        x = M + i * (cw + 0.16)
        _card(s, x, 1.42, cw, 1.24)
        _text(s, x + 0.18, 1.58, cw - 0.36, 0.42, big, size=21, color=INK,
              font=HEAD_FONT, bold=True)
        _text(s, x + 0.18, 2.03, cw - 0.36, 0.24, label, size=10, color=INK2)
        _text(s, x + 0.18, 2.28, cw - 0.36, 0.3, sub, size=8.5, color=FAINT)

    points = _ai_list(narr, "summary_points", 3)
    if not points:
        # Without the model the slide still has to say something true, so it
        # falls back to the three biggest categories stated plainly.
        points = [{"title": f"№{i + 1} — {c['label']} ({c['code']}): {num(c['minutes'])} min",
                   "body": f"{c['events']} hodisa, jami yo'qotishning {pct(c['share'])} qismi."}
                  for i, c in enumerate(d["categories"][:3])]

    y = 2.92
    for i, p in enumerate(points):
        h = 0.92
        _card(s, M, y, CW, h)
        _rect(s, M, y, 0.055, h, fill=GOLD)
        _text(s, M + 0.28, y + 0.16, CW - 0.56, 0.28,
              _sentence(p.get("title", "")), size=12.5, color=INK,
              font=HEAD_FONT, bold=True)
        _text(s, M + 0.28, y + 0.47, CW - 0.56, 0.4,
              p.get("body", ""), size=10, color=MUTED, line=1.18)
        y += h + 0.14

    head = _ai(narr, "summary_headline")
    _text(s, M, y + 0.04, CW, 0.46, head or NO_AI,
          size=10.5, color=INK2 if head else FAINT, bold=bool(head), line=1.2)


# ── slide 3 · what each category means, and what it cost ─────────────────────
def _glossary(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "Lug'at", "Kutish toifalari va haftalik natija", page, footer)

    # Canonical A→Z, not by size: this is the reference slide, and a reader
    # looking up «D3» should find it where the alphabet puts it.
    order = {name: i for i, (name, *_r) in enumerate(CATS)}
    rows = sorted(d["categories"], key=lambda c: order.get(c["name"], 99))

    col_w = (CW - 0.24) / 2
    rh, top = 0.53, 1.38
    per_col = (len(rows) + 1) // 2
    for i, c in enumerate(rows):
        col, idx = i // per_col, i % per_col
        x = M + col * (col_w + 0.24)
        y = top + idx * (rh + 0.07)
        _card(s, x, y, col_w, rh)
        _chip(s, x + 0.12, y + 0.11, 0.46, 0.31, c["code"],
              fill=c["colour"], color=WHITE, size=10.5)
        # The four boxes on a glossary row are stacked in two columns and must
        # not overlap even by a hundredth of an inch — see `_no_overlap` in
        # `self_check`: a box that overlaps its neighbour is a box whose text
        # will overlap the moment the text grows.
        _text(s, x + 0.68, y + 0.09, col_w - 2.0, 0.21, c["label"],
              size=10.5, color=INK, bold=True)
        note = c["full"] if c["full"] != c["label"] else ""
        _text(s, x + 0.68, y + 0.3, col_w - 2.0, 0.2, note, size=8, color=FAINT)
        _text(s, x + col_w - 1.26, y + 0.1, 1.14, 0.2,
              f"{num(c['minutes'])} min", size=11, color=INK, bold=True,
              align=PP_ALIGN.RIGHT, font=HEAD_FONT)
        _text(s, x + col_w - 1.26, y + 0.31, 1.14, 0.2,
              f"{c['events']} hodisa", size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)

    _text(s, M, H - 0.96, CW, 0.34,
          f"Jami {num(d['total'])} daqiqa · {d['events']} hodisa. "
          f"Ro'yxatda shu hafta qayd etilgan toifalar keltirilgan.",
          size=9, color=FAINT, line=1.2)


# ── slide 4 · which category took the time ───────────────────────────────────
def _categories(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "Toifalar kesimida", "Qaysi toifa qancha vaqt oldi?", page, footer)

    top = [c for c in d["categories"] if c["minutes"] > 0][:10]
    if top:
        # Ascending, because a horizontal bar chart draws its first category at
        # the BOTTOM — the biggest bar has to end up on top.
        rows = list(reversed(top))
        _bar_chart(s, M - 0.1, 1.35, 7.9, 4.9,
                   [f"{c['code']} · {c['label']}" for c in rows],
                   [c["minutes"] for c in rows],
                   [c["colour"] for c in rows])

    x = M + 8.05
    cw = CW - 8.05
    three = d["categories"][:3]
    share = sum(c["share"] for c in three)
    _card(s, x, 1.35, cw, 1.12, fill=INK)
    _text(s, x + 0.22, 1.5, cw - 0.44, 0.5, pct(share), size=30, color=GOLD2,
          font=HEAD_FONT, bold=True)
    _text(s, x + 0.22, 2.02, cw - 0.44, 0.36,
          "jami yo'qotish eng katta uchta toifada",
          size=9.5, color=ONDARK, line=1.15)

    y = 2.62
    for c in three:
        _card(s, x, y, cw, 0.86)
        _rect(s, x, y, 0.055, 0.86, fill=c["colour"])
        _chip(s, x + 0.2, y + 0.14, 0.44, 0.28, c["code"],
              fill=c["colour"], color=WHITE, size=9.5)
        _text(s, x + 0.74, y + 0.13, cw - 0.94, 0.24, c["label"],
              size=10.5, color=INK, bold=True)
        _text(s, x + 0.74, y + 0.37, cw - 0.94, 0.22,
              f"{num(c['minutes'])} min · {c['events']} hodisa · {pct(c['share'])}",
              size=9, color=MUTED)
        _text(s, x + 0.2, y + 0.6, cw - 0.4, 0.2, c["full"], size=8, color=FAINT)
        y += 0.94

    note = _ai(narr, "others_note")
    _text(s, M, 6.42, CW, 0.6, note or NO_AI,
          size=9.5, color=MUTED if note else FAINT, line=1.2)


# ── slide 5 · against last week ──────────────────────────────────────────────
def _comparison(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "Taqqoslash", "O'tgan hafta bilan: nima o'zgardi?", page, footer)

    _text(s, M, 1.3, CW, 0.24,
          f"{d['prev_period']}  →  {d['period']}", size=9.5, color=FAINT)

    better = _ai(narr, "compare_better")
    worse = _ai(narr, "compare_worse")
    cw = (CW - 0.2) / 2
    for i, (title, body, tone, bg) in enumerate([
        (f"Yaxshilandi — jami {signed_pct(d['delta_pct'])}" if d["delta_pct"] < 0
         else "Yaxshilangan yo'nalishlar", better, GREEN, GREENBG),
        ("Diqqat talab qiladi", worse, RED, REDBG),
    ]):
        x = M + i * (cw + 0.2)
        _card(s, x, 1.6, cw, 1.5)
        _rect(s, x, 1.6, cw, 0.36, fill=bg)
        _text(s, x + 0.2, 1.67, cw - 0.4, 0.24, title, size=10.5, color=tone, bold=True)
        _text(s, x + 0.2, 2.08, cw - 0.4, 0.9, body or NO_AI,
              size=9.5, color=INK2 if body else FAINT, line=1.22)

    # Every category, both weeks, biggest mover first — the table is what makes
    # the two sentences above checkable.
    rows = sorted((c for c in d["categories"] if c["minutes"] or c["prev"]),
                  key=lambda c: -abs(c["minutes"] - c["prev"]))[:9]
    y = 3.32
    hdr = [("Toifa", 3.3, PP_ALIGN.LEFT), ("O'tgan", 1.5, PP_ALIGN.RIGHT),
           ("Bu hafta", 1.5, PP_ALIGN.RIGHT), ("O'zgarish", 1.7, PP_ALIGN.RIGHT)]
    _rect(s, M, y, CW, 0.34, fill=INNER, line=LINE, radius=0.06)
    x = M + 0.2
    for label, w, al in hdr:
        _text(s, x, y + 0.09, w, 0.2, label, size=8.5, color=MUTED, bold=True,
              align=al, caps=True)
        x += w + 0.15
    y += 0.4

    for c in rows:
        delta = c["minutes"] - c["prev"]
        tone = RED if delta > 0 else (GREEN if delta < 0 else FAINT)
        x = M + 0.2
        _chip(s, x, y + 0.02, 0.42, 0.26, c["code"], fill=c["colour"], color=WHITE, size=8.5)
        _text(s, x + 0.5, y + 0.04, 2.7, 0.22, c["label"], size=9.5, color=INK)
        x += hdr[0][1] + 0.15
        _text(s, x, y + 0.04, hdr[1][1], 0.22, num(c["prev"]), size=9.5,
              color=MUTED, align=PP_ALIGN.RIGHT)
        x += hdr[1][1] + 0.15
        _text(s, x, y + 0.04, hdr[2][1], 0.22, num(c["minutes"]), size=9.5,
              color=INK, bold=True, align=PP_ALIGN.RIGHT)
        x += hdr[2][1] + 0.15
        arrow = "▲" if delta > 0 else ("▼" if delta < 0 else "—")
        _text(s, x, y + 0.04, hdr[3][1], 0.22,
              f"{arrow} {num(abs(delta))}" if delta else "—",
              size=9.5, color=tone, bold=True, align=PP_ALIGN.RIGHT)
        _rect(s, M + 0.2, y + 0.31, CW - 0.4, 0.008, fill=LINE)
        y += 0.36

    # No note here about where either week's minutes came from — see the
    # ruling recorded in `collect()`.


def cat_key(raw: str) -> str:
    """«Cat D3» / «D3» / «D3 — Otdellararo mahsulot» → «D3».

    The model is asked for a category code and answers with whatever reads
    naturally to it — sometimes the bare code, sometimes the code with its
    label attached, and the two arrive in the SAME response from different
    fields. A lookup that only accepted one spelling silently dropped the
    other, so a slide lost its root-cause line for no reason a reader could
    see. Everything after the first separator is discarded, and the result is
    matched case-insensitively against the canonical codes.
    """
    t = (raw or "").strip()
    if not t:
        return ""
    if t.lower().startswith("cat "):
        t = t[4:]
    for sep in ("—", "–", "-", ":", "·", "("):
        if sep in t:
            t = t.split(sep, 1)[0]
    t = t.strip().upper()
    codes = {c[1].upper(): c[1] for c in CATS}
    return codes.get(t, t)


def _top_events(events: list[dict], n: int) -> list[dict]:
    """The n longest waits in a group. Longest, not most recent: the reader is
    being shown where the time went."""
    return sorted(events, key=lambda e: -e["minutes"])[:n]


# An event's note is the one string on this deck whose length nobody here
# controls — a leader types it into `/idle-cell`. It gets up to EV_NOTE_LINES
# lines, and the meta line beneath it is placed at whatever height the note
# actually took. That is the fix for the overlap of 2026-09-04: the note used
# to be capped at 96 CHARACTERS inside a box one line tall, with the meta line
# nailed 0.21" below it, so any note that ran to two lines was drawn straight
# through it. A character cap cannot express «one line» — the card is a third
# of the page on one slide and a quarter of it on another.
EV_NOTE_LINES = 2
EV_NOTE_SIZE = 9
EV_META_SIZE = 8
EV_GAP = 0.09          # the air between one event and the next


def _event_block(e: dict, w: float) -> tuple[list[str], str, float]:
    """(note lines, meta line, total height) for one event, measured before
    anything is drawn — a stack has to know what an item costs to know whether
    it still fits."""
    tw = w - 0.16
    note, _ = deck_text.fit(f"«{_sentence(latin(e['note']))}»", tw, EV_NOTE_SIZE,
                            BODY_FONT, max_lines=EV_NOTE_LINES)
    mark = "" if e["stopped"] else " · to'xtatmagan"
    meta = (f"{num(e['minutes'])} min · {e['cell']} · {day_label(e['date_obj'])}"
            f" {e['start']}–{e['end']}{mark}")
    h = (deck_text.block_h_in(len(note), EV_NOTE_SIZE)
         + deck_text.block_h_in(1, EV_META_SIZE) + EV_GAP)
    return note, meta, h


def _events(slide, x, y: float, w: float, events: list[dict], *, bottom: float,
            colour=None, colour_fn=None) -> float:
    """A stack of events drawn against a vertical BUDGET, returning where it
    ended.

    `bottom` is the first inch this stack may not touch — the card's own edge,
    or whatever is drawn under it. Events are laid out while they fit and the
    stack stops when the next one would cross it. The step used to be a
    constant 0.5", which is only correct while every note happens to be one
    line: two long notes in a row walked the last event off its card.
    """
    for e in events:
        note, meta, h = _event_block(e, w)
        if y + h - EV_GAP > bottom:
            break
        nh = deck_text.block_h_in(len(note), EV_NOTE_SIZE)
        _rect(slide, x, y, 0.035, h - EV_GAP, fill=(colour_fn(e) if colour_fn
                                                    else colour) or GOLD)
        _text(slide, x + 0.16, y, w - 0.16, nh, "\n".join(note),
              size=EV_NOTE_SIZE, color=INK2, max_lines=len(note))
        _text(slide, x + 0.16, y + nh, w - 0.16,
              deck_text.block_h_in(1, EV_META_SIZE), meta,
              size=EV_META_SIZE, color=FAINT)
        y += h
    return y


# ── slide 6 · the three units the week cost most ─────────────────────────────
def _top_supervisors(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "So'ralgan kesim",
            "Top-3 brigadir: eng katta kutishlar kimda va nimadan?", page, footer)

    pains = {p.get("name"): p.get("pain") for p in _ai_list(narr, "supervisor_pains", 6)}
    three = d["supervisors"][:3]
    cw = (CW - 2 * 0.18) / 3
    for i, sup in enumerate(three):
        x = M + i * (cw + 0.18)
        _card(s, x, 1.32, cw, 4.86)
        _rect(s, x, 1.32, cw, 0.055, fill=GOLD)

        _chip(s, x + 0.18, 1.5, 0.3, 0.3, str(i + 1), fill=INK, color=GOLD2, size=11)
        _text(s, x + 0.58, 1.5, cw - 0.76, 0.28,
              short_name(sup["name"], cw - 0.76, 13, font=HEAD_FONT, bold=True),
              size=13, color=INK, font=HEAD_FONT, bold=True)
        _text(s, x + 0.58, 1.8, cw - 0.76, 0.22,
              f"{num(sup['minutes'])} min · jamining {pct(sup['share'])}",
              size=9.5, color=BROWN, bold=True)

        own = d["by_sup_events"].get(sup["name"], [])
        codes = sorted({e["cell"] for e in own})
        _text(s, x + 0.18, 2.1, cw - 0.36, 0.24,
              _clip(" · ".join(codes), 46) if codes else "yacheyka qayd etilmagan",
              size=8.5, color=FAINT)

        _events(s, x + 0.18, 2.44, cw - 0.36, _top_events(own, 5),
                bottom=5.36, colour_fn=lambda e: cat_meta(e["category"])[3])

        pain = pains.get(sup["name"])
        _rect(s, x + 0.18, 5.42, cw - 0.36, 0.72, fill=INNER, radius=0.06)
        _text(s, x + 0.32, 5.5, cw - 0.64, 0.58,
              pain or (f"{sup['events']} hodisa · eng ko'p: "
                       f"{cat_meta(sup['top_cat'])[1] if sup['top_cat'] else '—'}"),
              size=8.5, color=INK2 if pain else FAINT, line=1.18)

    rest = d["supervisors"][3:]
    if rest:
        _text(s, M, 6.35, CW, 0.4,
              "Qolgan brigadirlar: " + " · ".join(
                  f"{r['name']} {num(r['minutes'])}" for r in rest[:8]),
              size=8.5, color=FAINT, line=1.2)


# ── slide 7 · the three categories, opened up ────────────────────────────────
def _top_categories(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "So'ralgan kesim", "Top-3 toifa: aynan nima kutildi?", page, footer)

    roots = {cat_key(r.get("cat")): r.get("root")
             for r in _ai_list(narr, "category_roots", 6)}
    three = d["categories"][:3]
    cw = (CW - 2 * 0.18) / 3
    for i, c in enumerate(three):
        x = M + i * (cw + 0.18)
        _card(s, x, 1.32, cw, 4.86)
        _rect(s, x, 1.32, cw, 0.055, fill=c["colour"])

        _chip(s, x + 0.18, 1.5, 0.46, 0.3, c["code"], fill=c["colour"], color=WHITE, size=11)
        _text(s, x + 0.74, 1.5, cw - 0.92, 0.28, c["label"],
              size=12.5, color=INK, font=HEAD_FONT, bold=True)
        _text(s, x + 0.74, 1.8, cw - 0.92, 0.22,
              f"{num(c['minutes'])} min · {c['events']} hodisa · {pct(c['share'])}",
              size=9.5, color=BROWN, bold=True)

        own = d["by_cat_events"].get(c["name"], [])
        who = sorted({e["supervisor"] for e in own})
        _text(s, x + 0.18, 2.1, cw - 0.36, 0.24,
              _clip(" · ".join(who), 46) if who else "—", size=8.5, color=FAINT)

        _events(s, x + 0.18, 2.44, cw - 0.36, _top_events(own, 5),
                bottom=5.36, colour=c["colour"])

        root = roots.get(c["code"])
        _rect(s, x + 0.18, 5.42, cw - 0.36, 0.72, fill=INNER, radius=0.06)
        _text(s, x + 0.32, 5.5, cw - 0.64, 0.58,
              ("Ildiz: " + root) if root else c["full"],
              size=8.5, color=INK2 if root else FAINT, line=1.18)

    _text(s, M, 6.35, CW, 0.4,
          f"Uchtasi birgalikda {num(sum(c['minutes'] for c in three))} daqiqa — "
          f"jami yo'qotishning {pct(sum(c['share'] for c in three))} qismi.",
          size=9.5, color=MUTED)


# ── slide 8 · everything else ────────────────────────────────────────────────
def _other_categories(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "Qolgan toifalar", "Uchlikdan tashqarida nima bo'ldi?", page, footer)

    rest = [c for c in d["categories"][3:] if c["minutes"] > 0][:3]
    if not rest:
        _text(s, M, 3.2, CW, 0.5,
              "Bu hafta boshqa toifalarda kutish qayd etilmagan.",
              size=12, color=FAINT, align=PP_ALIGN.CENTER)
        return

    cw = (CW - (len(rest) - 1) * 0.18) / len(rest)
    for i, c in enumerate(rest):
        x = M + i * (cw + 0.18)
        _card(s, x, 1.32, cw, 4.0)
        _rect(s, x, 1.32, cw, 0.055, fill=c["colour"])
        _chip(s, x + 0.18, 1.5, 0.46, 0.3, c["code"], fill=c["colour"], color=WHITE, size=11)
        _text(s, x + 0.74, 1.5, cw - 0.92, 0.28, c["label"],
              size=12.5, color=INK, font=HEAD_FONT, bold=True)
        delta = c["minutes"] - c["prev"]
        move = f" ({signed_pct(delta / c['prev'] * 100)})" if c["prev"] else ""
        _text(s, x + 0.74, 1.8, cw - 0.92, 0.22,
              f"{num(c['minutes'])} min · {c['events']} hodisa{move}",
              size=9.5, color=BROWN, bold=True)
        _events(s, x + 0.18, 2.2, cw - 0.36,
                _top_events(d["by_cat_events"].get(c["name"], []), 5),
                bottom=5.18, colour=c["colour"])

    note = _ai(narr, "others_note")
    _rect(s, M, 5.52, CW, 0.86, fill=INNER, line=LINE, radius=0.08)
    _text(s, M + 0.24, 5.66, CW - 0.48, 0.62, note or NO_AI,
          size=9.5, color=INK2 if note else FAINT, line=1.2)


# ── slide 9 · every unit, ranked ─────────────────────────────────────────────
def _supervisors(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    three = d["supervisors"][:3]
    share = sum(x["share"] for x in three)
    _chrome(s, "Brigadirlar kesimida",
            f"Top-3 brigadir jamining {pct(share)} ini beradi", page, footer)

    rows = [x for x in d["supervisors"] if x["minutes"] > 0][:12]
    if rows:
        rev = list(reversed(rows))
        # Red where the unit crossed the 50-minute day flag at least once —
        # the same threshold /downtime paints its own table with.
        colours = [RED if r["flagged_days"] else RGBColor(0x63, 0x66, 0xF1) for r in rev]
        _bar_chart(s, M - 0.1, 1.3, 8.15, 5.0,
                   [_clip(r["name"], 26) for r in rev],
                   [r["minutes"] for r in rev], colours)

    x = M + 8.3
    cw = CW - 8.3
    _text(s, x, 1.32, cw, 0.24, "IZOH", size=9, color=GOLD, bold=True)
    y = 1.62
    for i, r in enumerate(three):
        _card(s, x, y, cw, 0.78)
        _text(s, x + 0.18, y + 0.12, cw - 0.36, 0.24,
              f"{i + 1}. {short_name(r['name'], cw - 0.72, 10.5, bold=True)}",
              size=10.5, color=INK, bold=True)
        top = cat_meta(r["top_cat"])[1] if r["top_cat"] else "—"
        _text(s, x + 0.18, y + 0.38, cw - 0.36, 0.36,
              f"{num(r['minutes'])} min · {r['events']} hodisa\neng ko'p: {top}",
              size=8.5, color=MUTED, line=1.18)
        y += 0.86

    _rect(s, x, y + 0.1, cw, 1.5, fill=INNER, line=LINE, radius=0.08)
    _text(s, x + 0.18, y + 0.24, cw - 0.36, 1.25,
          "Katta raqam — yomon ish degani emas. Bu yerda brigadir KUTGAN "
          "vaqt ko'rsatilgan, ya'ni unga boshqa bo'limlardan yetib kelmagan "
          "narsa. Qizil ustun — hech bo'lmaganda bir kun 50 daqiqadan oshgan "
          "brigadir.",
          size=8.5, color=MUTED, line=1.22)

    if d["quality"]["silent_supervisors"]:
        _text(s, M, H - 0.86, CW, 0.28,
              "Jurnal bermagan brigadirlar: "
              + _clip(" · ".join(d["quality"]["silent_supervisors"]), 150),
              size=8.5, color=FAINT)


# ── slide 10 · the week, day by day ──────────────────────────────────────────
def _daily(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    peak = max(d["daily"], key=lambda x: x["minutes"]) if d["daily"] else None
    title = (f"Cho'qqi — {day_full(peak['date'])}" if peak and peak["minutes"]
             else "Hafta dinamikasi")
    _chrome(s, "Hafta dinamikasi", title, page, footer)

    vals = [x["minutes"] for x in d["daily"]]
    top = max(vals) if vals else 0
    colours = [RED if v == top and v else RGBColor(0xC9, 0x9E, 0x63) for v in vals]
    _bar_chart(s, M - 0.1, 1.35, 8.6, 4.3,
               [x["label"] for x in d["daily"]], vals, colours,
               horizontal=False, gap=45)

    x = M + 8.75
    cw = CW - 8.75
    if peak and peak["minutes"]:
        _card(s, x, 1.35, cw, 1.5, fill=INK)
        _text(s, x + 0.22, 1.5, cw - 0.44, 0.5, f"{num(peak['minutes'])} min",
              size=26, color=GOLD2, font=HEAD_FONT, bold=True)
        _text(s, x + 0.22, 2.02, cw - 0.44, 0.7, day_full(peak["date"]),
              size=10, color=ONDARK, line=1.18)

        same = [e for e in d["events_on"] if e["date_obj"] == peak["date"]]
        _events(s, x, 3.0, cw, _top_events(same, 4), bottom=6.82,
                colour_fn=lambda e: cat_meta(e["category"])[3])

    note = _ai(narr, "daily_note")
    _rect(s, M, 5.86, 8.5, 0.82, fill=INNER, line=LINE, radius=0.08)
    _text(s, M + 0.22, 5.99, 8.1, 0.6, note or NO_AI,
          size=9.5, color=INK2 if note else FAINT, line=1.2)

    avg = (sum(vals) / len([v for v in vals if v])) if any(vals) else 0
    _text(s, M, H - 0.86, 8.5, 0.26,
          f"Kunlik o'rtacha (kutish bo'lgan kunlar bo'yicha): {num(avg)} daqiqa.",
          size=8.5, color=FAINT)


# ── slide 11 · which cells stopped most ──────────────────────────────────────
def _cells(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "Yacheykalar kesimida", "Eng ko'p to'xtagan yacheykalar", page, footer)

    rows = d["cells"][:12]
    if rows:
        rev = list(reversed(rows))
        _bar_chart(s, M - 0.1, 1.3, 8.15, 4.7,
                   [r["code"] for r in rev], [r["minutes"] for r in rev],
                   [RGBColor(0x14, 0xB8, 0xA6)] * len(rev))

    x = M + 8.3
    cw = CW - 8.3
    _text(s, x, 1.32, cw, 0.24, "ENG OG'IRLARI", size=9, color=GOLD, bold=True)
    y = 1.62
    for r in rows[:4]:
        _card(s, x, y, cw, 0.82)
        _text(s, x + 0.18, y + 0.12, cw - 0.36, 0.24,
              f"{r['code']} — {num(r['minutes'])} min", size=11, color=INK,
              font=HEAD_FONT, bold=True)
        who = " · ".join(v for v in (r.get("leader"), r.get("supervisor")) if v)
        _text(s, x + 0.18, y + 0.38, cw - 0.36, 0.34,
              f"{_one(who, cw - 0.36, 8.5)}\n{r['events']} hodisa",
              size=8.5, color=MUTED, line=1.18)
        y += 0.9

    note = _ai(narr, "cells_note")
    _rect(s, M, 6.06, 8.5, 0.62, fill=INNER, line=LINE, radius=0.08)
    _text(s, M + 0.22, 6.17, 8.1, 0.44, note or NO_AI,
          size=9.5, color=INK2 if note else FAINT, line=1.2)

    # The bar chart above is the cells' OWN arithmetic and does not add up to
    # the headline — saying so here is cheaper than a reader discovering it.
    _text(s, M, H - 0.86, CW, 0.28,
          "Yacheyka jamlanmasi — o'sha yacheykaning kunlik kutish oralig'i "
          "birlashmasi. Brigadir ko'rsatkichi esa yacheykalarning odam soniga "
          "vaznlangan o'rtachasi, shuning uchun bu ustunlar yig'indisi umumiy "
          "raqamdan katta chiqadi.", size=8.5, color=FAINT, line=1.2)


# ── slide 12 · what to do ────────────────────────────────────────────────────
def _actions(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "Tavsiyalar", "Nima qilish kerak?", page, footer)

    actions = _ai_list(narr, "actions", 4)
    if not actions:
        actions = [{"cat": c["code"],
                    "text": f"{c['label']} — {num(c['minutes'])} daqiqa, "
                            f"{c['events']} hodisa. Sabablarini brigadirlar bilan ko'rib chiqish."}
                   for c in d["categories"][:3]]

    by_code = {c["code"]: c for c in d["categories"]}
    y = 1.36
    for i, a in enumerate(actions):
        code = cat_key(a.get("cat"))
        c = by_code.get(code)
        h = 1.06
        _card(s, M, y, CW, h)
        _rect(s, M, y, 0.055, h, fill=c["colour"] if c else GOLD)
        _chip(s, M + 0.24, y + 0.2, 0.5, 0.3, code or str(i + 1),
              fill=c["colour"] if c else GOLD, color=WHITE, size=10.5)
        if c:
            _text(s, M + 0.86, y + 0.18, 3.2, 0.24, c["label"], size=10.5,
                  color=INK, bold=True)
            _text(s, M + 0.86, y + 0.42, 3.2, 0.22,
                  f"{num(c['minutes'])} min · {pct(c['share'])}",
                  size=9, color=BROWN, bold=True)
        _text(s, M + 4.25, y + 0.2, CW - 4.5, 0.72,
              a.get("text", ""), size=10, color=INK2, line=1.22)
        y += h + 0.13

    covered = sum(by_code[c]["minutes"]
                  for c in {cat_key(a.get("cat")) for a in actions} if c in by_code)
    if covered and d["total"]:
        _rect(s, M, y + 0.06, CW, 0.62, fill=INK, radius=0.08)
        _text(s, M + 0.24, y + 0.2, CW - 0.48, 0.36,
              f"Bu choralar {num(covered)} daqiqani — jami yo'qotishning "
              f"{pct(covered / d['total'] * 100)} qismini qamrab oladi.",
              size=10.5, color=GOLD2, bold=True)


# ── slide 13 · the conclusion ────────────────────────────────────────────────
def _conclusion(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=DARK)
    _rect(s, 0, 0, W, 0.09, fill=GOLD)
    _chrome(s, "Xulosa", "", page, footer, dark=True)

    head = _ai(narr, "conclusion_headline")
    _text(s, M, 1.1, CW - 1.0, 1.42,
          head or f"{hours(d['total'])} soat kutish — {d['events']} hodisa, "
                  f"{d['cell_count']} yacheykada.",
          size=26, color=WHITE, font=HEAD_FONT, bold=True, line=1.12)

    three = d["supervisors"][:3]
    tcat = d["categories"][:3]
    _text(s, M, 2.62, CW, 0.3,
          f"Top-3 brigadir — {pct(sum(x['share'] for x in three))}; "
          f"top-3 toifa ({' · '.join(c['code'] for c in tcat)}) — "
          f"{pct(sum(c['share'] for c in tcat))}.",
          size=11, color=ONDARK)

    points = _ai_list(narr, "conclusion_points", 3)
    if not points:
        points = [{"title": c["label"],
                   "body": f"{num(c['minutes'])} daqiqa, {c['events']} hodisa."}
                  for c in tcat]
    cw = (CW - 2 * 0.2) / 3
    for i, p in enumerate(points):
        x = M + i * (cw + 0.2)
        _rect(s, x, 3.0, cw, 2.5, fill=RGBColor(0x3A, 0x28, 0x1C), radius=0.12)
        _chip(s, x + 0.22, 3.2, 0.34, 0.34, str(i + 1), fill=GOLD2, color=DARK, size=12)
        _text(s, x + 0.22, 3.72, cw - 0.44, 0.5,
              _sentence(p.get("title", "")), size=13, color=WHITE,
              font=HEAD_FONT, bold=True, line=1.1)
        _text(s, x + 0.22, 4.3, cw - 0.44, 1.05,
              p.get("body", ""), size=9.5, color=ONDARK, line=1.22)

    nxt = d["window"][1] + timedelta(days=7)
    _text(s, M, 5.72, CW, 0.3,
          f"Keyingi hisobot: {day_label(nxt)} chorshanba yakunida.",
          size=10, color=GOLD2)


# ── slide 14 · how this was made, and where it is thin ───────────────────────
def _appendix(prs, d: dict, narr, page: int, footer: str):
    s = _blank(prs)
    _rect(s, 0, 0, W, H, fill=CREAM)
    _chrome(s, "Ilova", "Metodika va ma'lumot sifati", page, footer)

    q = d["quality"]
    cw = (CW - 0.2) / 2

    _card(s, M, 1.32, cw, 2.5)
    _text(s, M + 0.24, 1.48, cw - 0.48, 0.26, "Ma'lumot qanday olindi",
          size=11.5, color=INK, font=HEAD_FONT, bold=True)
    lines = [
        f"Manba: liderlarning «Yacheykalar kutishi» yozuvlari "
        f"(/idle-cell). Faqat {d['period']} oralig'i olindi.",
        f"To'xtagan hodisalar: {d['events']} ta, {num(d['total'])} daqiqa.",
        f"Brigadir ko'rsatkichi — yacheykalarning odam soniga vaznlangan "
        f"o'rtachasi: Σ(N·T)÷ΣN. Shu sababli yacheykalar yig'indisi undan katta.",
        f"Toifalar: yacheyka zagruzkasi hisobiga kiradigan kutish "
        f"toifalari olindi.",
    ]
    y = 1.82
    for t in lines:
        _rect(s, M + 0.24, y + 0.06, 0.055, 0.2, fill=GOLD)
        _text(s, M + 0.42, y, cw - 0.7, 0.46, t, size=9, color=MUTED, line=1.2)
        y += 0.5

    _card(s, M + cw + 0.2, 1.32, cw, 2.5)
    _text(s, M + cw + 0.44, 1.48, cw - 0.48, 0.26, "Hisobga KIRMAGANI",
          size=11.5, color=INK, font=HEAD_FONT, bold=True)
    out = [
        f"«To'xtatmagan» kutishlar: {d['events_ns']} ta, {num(d['total_ns'])} "
        f"daqiqa — yacheyka ishlashda davom etgan, shuning uchun asosiy "
        f"raqamga qo'shilmadi.",
    ]
    if q["silent_supervisors"]:
        out.append(f"Jurnal bermagan brigadirlar: {len(q['silent_supervisors'])} ta — "
                   f"{_clip(' · '.join(q['silent_supervisors']), 110)}.")
    y = 1.82
    for t in out:
        _rect(s, M + cw + 0.44, y + 0.06, 0.055, 0.2, fill=RED)
        _text(s, M + cw + 0.62, y, cw - 0.7, 0.6, t, size=9, color=MUTED, line=1.2)
        y += 0.62

    _card(s, M, 3.98, CW, 1.72)
    _text(s, M + 0.24, 4.14, CW - 0.48, 0.26, "Sifat bo'yicha kuzatuvlar",
          size=11.5, color=INK, font=HEAD_FONT, bold=True)
    facts = [
        (f"{q['thin_notes']} hodisada izoh juda qisqa "
         f"({THIN_NOTE_CHARS} belgidan kam) — sabab tiklanmaydi.",
         RED if q["thin_notes"] else GREEN),
        (f"{q['cyrillic_notes']} izoh kirill alifbosida yozilgan — bu faylda "
         f"lotinga o'girildi, matn o'zgartirilmadi.", MUTED),
        (f"{q['flag_days']} brigadir-kun 50 daqiqadan oshdi.",
         RED if q["flag_days"] else GREEN),
        (f"Taqqoslash davri: {d['prev_period']} ({num(d['prev_total'])} daqiqa). "
         f"Ikkala davr ham 8 kunlik va bir kunni baham ko'radi.", MUTED),
    ]
    y = 4.48
    for t, tone in facts:
        _rect(s, M + 0.24, y + 0.05, 0.055, 0.2, fill=tone)
        _text(s, M + 0.42, y, CW - 0.7, 0.3, t, size=9, color=MUTED, line=1.2)
        y += 0.32

    ai_line = ("Sahifalardagi izoh matnlari sun'iy intellekt tomonidan "
               "yozilgan; raqamlar, jadvallar va grafiklar to'g'ridan-to'g'ri "
               "ma'lumotlar bazasidan olingan. Liderlar yozuvlari qayta "
               "yozilmagan — faqat qo'shtirnoq ichida keltirilgan."
               if narr else
               "Bu faylda AI izohlari yo'q — barcha matn ma'lumotlar bazasidan "
               "olingan raqamlardan iborat.")
    _rect(s, M, 5.88, CW, 0.7, fill=INNER, line=LINE, radius=0.08)
    _text(s, M + 0.24, 6.0, CW - 0.48, 0.5, ai_line, size=8.5, color=FAINT, line=1.2)


# ── the deck ─────────────────────────────────────────────────────────────────
_SLIDES = [
    _summary, _glossary, _categories, _comparison, _top_supervisors,
    _top_categories, _other_categories, _supervisors, _daily, _cells,
    _actions, _conclusion, _appendix,
]


def build(d: dict, narrative: dict | None = None) -> bytes:
    """The finished .pptx. `narrative` may be None — see `deck_narrative`."""
    _TRIMS.clear()
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)

    footer = f"Yacheykalardagi kutish vaqtlari · {d['period']} · {d['factory']}"
    _cover(prs, d, narrative)
    for i, fn in enumerate(_SLIDES, start=2):
        fn(prs, d, narrative, i, footer)

    if _TRIMS:
        # Not a failure — the no-overflow rule doing its job — but the one
        # thing worth knowing about a finished file. A box that trims week
        # after week is a box that should be bigger.
        log.info("DECK built with %d trimmed text block(s)", len(_TRIMS))

    # The other half of the rule, checked on the real week rather than trusted:
    # `_text` keeps every string inside its own box, and this keeps the boxes
    # off each other. Findings are logged, never raised — an operator waiting
    # on Wednesday's report is not served by being handed nothing.
    for problem in deck_text.check_layout(prs, W, H):
        log.warning("DECK layout: %s", problem)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def filename(d: dict) -> str:
    a, b = d["window"]
    return (f"Yacheykalar_kutish_tahlili_{a.strftime('%d-%m')}_"
            f"{b.strftime('%d-%m-%Y')}.pptx")
